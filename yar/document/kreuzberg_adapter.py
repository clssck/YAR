"""Kreuzberg adapter for document parsing and chunking.

Kreuzberg is a polyglot document intelligence framework with Rust core,
supporting 56+ document formats with built-in semantic chunking for RAG.

This adapter provides a clean interface for YAR to use Kreuzberg
for document extraction and optional semantic chunking.
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import TYPE_CHECKING, Any

import tiktoken

from yar.utils import logger

# Cache the tokenizer for performance
_TOKENIZER_CACHE: dict[str, tiktoken.Encoding] = {}


def _get_tokenizer(model: str = 'gpt-4o') -> tiktoken.Encoding:
    """Get cached tiktoken tokenizer."""
    if model not in _TOKENIZER_CACHE:
        try:
            _TOKENIZER_CACHE[model] = tiktoken.encoding_for_model(model)
        except KeyError:
            # Fallback to cl100k_base for unknown models
            _TOKENIZER_CACHE[model] = tiktoken.get_encoding('cl100k_base')
    return _TOKENIZER_CACHE[model]


def tokens_to_chars(token_count: int, model: str = 'gpt-4o') -> int:
    """Convert token count to approximate character count using actual tokenizer.

    Uses a sample of common text to estimate the average chars/token ratio,
    then applies it to the token count. This is more accurate than the
    naive 4 chars/token approximation.

    Args:
        token_count: Number of tokens
        model: Model name for tokenizer selection

    Returns:
        Estimated character count
    """
    tokenizer = _get_tokenizer(model)

    # Mixed content sample (technical terms, abbreviations, numbers) for realistic ratio
    sample = (
        'CMC Cross-Sharing Session 2023: Key findings from Q3 regulatory submissions. '
        'FDA approved 3 NDAs for SAR-123456. EMA issued GMP warnings for Site-A. '
        'Process B vs Process C comparability data shows 95.2% similarity. '
        'def extract_entities(doc: Document) -> list[Entity]: return parser.parse(doc)'
    )

    sample_tokens = len(tokenizer.encode(sample))
    chars_per_token = len(sample) / sample_tokens if sample_tokens > 0 else 4.0

    return int(token_count * chars_per_token)


if TYPE_CHECKING:
    pass


def _setup_pdfium_for_kreuzberg() -> bool:
    """Set up pdfium library for Kreuzberg PDF support.

    Kreuzberg's RC version may not bundle pdfium on all platforms.
    This function creates a symlink from pypdfium2's bundled library
    to where kreuzberg expects it.

    Returns:
        bool: True if setup succeeded or was unnecessary, False otherwise
    """
    try:
        import kreuzberg
        import pypdfium2_raw

        kreuzberg_file = kreuzberg.__file__
        if not kreuzberg_file:
            logger.debug('kreuzberg.__file__ is unavailable; skipping pdfium setup')
            return False
        kreuzberg_dir = Path(kreuzberg_file).parent
        pdfium_target = kreuzberg_dir / 'libpdfium.dylib'

        # Skip if already exists
        if pdfium_target.exists():
            return True

        # Find pypdfium2's library
        pypdfium_file = pypdfium2_raw.__file__
        if not pypdfium_file:
            logger.debug('pypdfium2_raw.__file__ is unavailable; skipping pdfium setup')
            return False
        pypdfium_dir = Path(pypdfium_file).parent
        pdfium_source = pypdfium_dir / 'libpdfium.dylib'

        if not pdfium_source.exists():
            # Try Linux naming
            pdfium_source = pypdfium_dir / 'libpdfium.so'
            pdfium_target = kreuzberg_dir / 'libpdfium.so'

        if pdfium_source.exists() and not pdfium_target.exists():
            try:
                os.symlink(pdfium_source, pdfium_target)
                logger.debug(f'Created pdfium symlink: {pdfium_target} -> {pdfium_source}')
                return True
            except OSError as e:
                # May fail in read-only environments (Docker), which is fine
                # as Docker images should have this pre-configured
                logger.debug(f'Could not create pdfium symlink: {e}')
                return False

        return True
    except ImportError:
        # pypdfium2 not installed
        return False
    except Exception as e:
        logger.debug(f'pdfium setup error: {e}')
        return False


@lru_cache(maxsize=1)
def is_kreuzberg_available() -> bool:
    """Check if kreuzberg is available (cached check).

    This function uses lru_cache to avoid repeated import attempts.
    The result is cached after the first call.

    Returns:
        bool: True if kreuzberg is available, False otherwise
    """
    try:
        import kreuzberg  # noqa: F401

        return True
    except ImportError:
        return False


# Module-level constant for quick checks
KREUZBERG_AVAILABLE = is_kreuzberg_available()

# Auto-setup pdfium for PDF support (creates symlink from pypdfium2 if needed)
if KREUZBERG_AVAILABLE:
    _setup_pdfium_for_kreuzberg()


def _extract_ppsx_from_zip(file_path: str | Path) -> str:
    """Extract text from PPSX using direct ZIP extraction.

    PPSX files are structurally identical to PPTX (OpenXML package),
    but python-pptx rejects them due to content-type validation.
    This function extracts text directly from the slide XML files.

    Args:
        file_path: Path to the PPSX file

    Returns:
        Extracted text content
    """
    import re
    import zipfile

    all_text = []

    with zipfile.ZipFile(str(file_path), 'r') as zf:
        # Find all slide XML files
        slide_files = sorted(
            [n for n in zf.namelist() if re.match(r'ppt/slides/slide\d+\.xml', n)],
            key=lambda x: int(re.search(r'slide(\d+)', x).group(1)),  # type: ignore[union-attr]
        )

        for slide_num, slide_file in enumerate(slide_files, 1):
            slide_texts = [f'[Slide {slide_num}]']

            content = zf.read(slide_file).decode('utf-8')
            # Extract all <a:t> elements (text runs in DrawingML)
            texts = re.findall(r'<a:t[^>]*>([^<]*)</a:t>', content)
            for text in texts:
                text = text.strip()
                if text:
                    slide_texts.append(text)

            if len(slide_texts) > 1:
                all_text.extend(slide_texts)
                all_text.append('')

    return '\n'.join(all_text)


def _extract_ppsx_bytes_from_zip(data: bytes) -> str:
    """Extract text from PPSX bytes using direct ZIP extraction.

    Args:
        data: PPSX file content as bytes

    Returns:
        Extracted text content
    """
    import io
    import re
    import zipfile

    all_text = []

    with zipfile.ZipFile(io.BytesIO(data), 'r') as zf:
        slide_files = sorted(
            [n for n in zf.namelist() if re.match(r'ppt/slides/slide\d+\.xml', n)],
            key=lambda x: int(re.search(r'slide(\d+)', x).group(1)),  # type: ignore[union-attr]
        )

        for slide_num, slide_file in enumerate(slide_files, 1):
            slide_texts = [f'[Slide {slide_num}]']

            content = zf.read(slide_file).decode('utf-8')
            texts = re.findall(r'<a:t[^>]*>([^<]*)</a:t>', content)
            for text in texts:
                text = text.strip()
                if text:
                    slide_texts.append(text)

            if len(slide_texts) > 1:
                all_text.extend(slide_texts)
                all_text.append('')

    return '\n'.join(all_text)


def _extract_pptx_with_python_pptx(file_path: str | Path) -> str:
    """Fallback PPTX extraction using python-pptx.

    Used when Kreuzberg's native PPTX parser fails (e.g., on shapes without txBody).
    python-pptx handles missing text gracefully with hasattr checks.

    Args:
        file_path: Path to the PPTX file

    Returns:
        Extracted text content
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def extract_text_from_shape(shape: Any) -> list[str]:
        """Recursively extract text from a shape, handling groups."""
        texts = []
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                texts.extend(extract_text_from_shape(child_shape))
        elif hasattr(shape, 'text') and shape.text:
            texts.append(shape.text.strip())
        # Also check for tables - must verify shape actually has a table
        if shape.has_table:
            try:
                for row in shape.table.rows:
                    row_texts = []
                    for cell in row.cells:
                        if cell.text:
                            row_texts.append(cell.text.strip())
                    if row_texts:
                        texts.append(' | '.join(row_texts))
            except Exception:
                pass  # Shape reports has_table but table access fails
        return texts

    prs = Presentation(str(file_path))
    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f'[Slide {slide_num}]']
        for shape in slide.shapes:
            slide_texts.extend(extract_text_from_shape(shape))
        if len(slide_texts) > 1:  # More than just the slide marker
            all_text.extend(slide_texts)
            all_text.append('')  # Blank line between slides

    return '\n'.join(all_text)


def _extract_pptx_bytes_with_python_pptx(data: bytes) -> str:
    """Fallback PPTX extraction from bytes using python-pptx.

    Args:
        data: PPTX file content as bytes

    Returns:
        Extracted text content
    """
    import io

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def extract_text_from_shape(shape: Any) -> list[str]:
        """Recursively extract text from a shape, handling groups."""
        texts = []
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child_shape in shape.shapes:
                texts.extend(extract_text_from_shape(child_shape))
        elif hasattr(shape, 'text') and shape.text:
            texts.append(shape.text.strip())
        # Also check for tables - must verify shape actually has a table
        if shape.has_table:
            try:
                for row in shape.table.rows:
                    row_texts = []
                    for cell in row.cells:
                        if cell.text:
                            row_texts.append(cell.text.strip())
                    if row_texts:
                        texts.append(' | '.join(row_texts))
            except Exception:
                pass  # Shape reports has_table but table access fails
        return texts

    prs = Presentation(io.BytesIO(data))
    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f'[Slide {slide_num}]']
        for shape in slide.shapes:
            slide_texts.extend(extract_text_from_shape(shape))
        if len(slide_texts) > 1:  # More than just the slide marker
            all_text.extend(slide_texts)
            all_text.append('')  # Blank line between slides

    return '\n'.join(all_text)


@dataclass
class ChunkingOptions:
    """Kreuzberg chunking config."""

    enabled: bool = False
    max_chars: int = 4170  # ~1000 tokens via tokens_to_chars
    max_overlap: int = 417  # ~100 tokens
    preset: str | None = 'semantic'  # 'recursive', 'semantic', or None


@dataclass
class OcrOptions:
    """Kreuzberg OCR config."""

    backend: str = 'tesseract'
    language: str = 'en'
    enable_table_detection: bool = True


@dataclass
class PdfOptions:
    """Kreuzberg PDF-specific config."""

    extract_images: bool = False
    extract_metadata: bool = True
    enable_hierarchy: bool = False  # Document structure analysis


@dataclass
class LanguageDetectionOptions:
    """Kreuzberg language detection config."""

    enabled: bool = True
    min_confidence: float = 0.5
    detect_multiple: bool = False


@dataclass
class TokenReductionOptions:
    """Kreuzberg token reduction config for LLM optimization."""

    mode: str | None = None  # 'off', 'light', 'moderate', 'aggressive', 'maximum'
    preserve_important_words: bool = True


@dataclass
class PageOptions:
    """Kreuzberg page tracking config."""

    extract_pages: bool = True
    insert_page_markers: bool = True  # Adds [Page N] markers for source citation
    marker_format: str | None = None  # Custom format, e.g., "--- Page {page} ---"


@dataclass
class HierarchyOptions:
    """Kreuzberg document structure detection via font clustering."""

    enabled: bool = False
    k_clusters: int = 6  # Number of heading levels to detect
    include_bbox: bool = False  # Include bounding box coordinates
    ocr_coverage_threshold: float | None = None


@dataclass
class ExtractionOptions:
    """Combined extraction options for Kreuzberg.

    New in kreuzberg 4.2:
        output_format: Convert extracted content to plain/markdown/html/djot
        result_format: 'unified' (default) or 'element_based' for semantic elements
    """

    chunking: ChunkingOptions | None = None
    ocr: OcrOptions | None = None
    pdf: PdfOptions | None = None
    language_detection: LanguageDetectionOptions | None = None
    token_reduction: TokenReductionOptions | None = None
    pages: PageOptions | None = None
    hierarchy: HierarchyOptions | None = None
    mime_type: str | None = None
    use_cache: bool = True
    enable_quality_processing: bool = True
    force_ocr: bool = False  # Force OCR even for text-based documents (fallback for PPTX parsing failures)
    # New in kreuzberg 4.2.0
    output_format: str | None = None  # 'plain', 'markdown', 'html', 'djot' - converts ANY file to this format
    result_format: str | None = None  # 'unified' (default) or 'element_based' for semantic extraction


@dataclass
class TextChunk:
    """A chunk of extracted text with metadata.

    Attributes:
        content: The text content of the chunk
        index: Zero-based chunk index
        start_char: Starting character offset in original document
        end_char: Ending character offset in original document
        metadata: Additional metadata from Kreuzberg
    """

    content: str
    index: int
    start_char: int | None = None
    end_char: int | None = None
    metadata: dict[str, Any] | None = None


@dataclass
class ExtractionResult:
    """Result of document extraction.

    Attributes:
        content: Full extracted text content
        chunks: List of text chunks (if chunking enabled)
        mime_type: Detected or specified MIME type
        metadata: Document metadata
        tables: Extracted tables (if any)
        detected_languages: Languages detected in document
    """

    content: str
    chunks: list[TextChunk] | None = None
    mime_type: str | None = None
    metadata: dict[str, Any] | None = None
    tables: list[Any] | None = None
    detected_languages: list[str] | None = None


def _build_extraction_config(options: ExtractionOptions | None = None) -> Any:
    """Build Kreuzberg ExtractionConfig from our options."""
    from kreuzberg import (
        ChunkingConfig,
        ExtractionConfig,
        LanguageDetectionConfig,
        OcrConfig,
        PdfConfig,
        TokenReductionConfig,
    )

    config_kwargs: dict[str, Any] = {}

    if not options:
        return None

    if options.use_cache is not None:
        config_kwargs['use_cache'] = options.use_cache

    if options.enable_quality_processing is not None:
        config_kwargs['enable_quality_processing'] = options.enable_quality_processing

    if options.force_ocr:
        config_kwargs['force_ocr'] = True

    if options.chunking and options.chunking.enabled:
        chunking_kwargs: dict[str, Any] = {
            'max_chars': options.chunking.max_chars,
            'max_overlap': options.chunking.max_overlap,
        }
        if options.chunking.preset:
            chunking_kwargs['preset'] = options.chunking.preset
        config_kwargs['chunking'] = ChunkingConfig(**chunking_kwargs)

    if options.ocr:
        ocr_kwargs: dict[str, Any] = {'backend': options.ocr.backend}
        if options.ocr.language:
            ocr_kwargs['language'] = options.ocr.language
        if options.ocr.backend == 'tesseract' and options.ocr.enable_table_detection:
            from kreuzberg import TesseractConfig

            ocr_kwargs['tesseract_config'] = TesseractConfig(enable_table_detection=True)
        config_kwargs['ocr'] = OcrConfig(**ocr_kwargs)

    if options.pdf:
        pdf_kwargs: dict[str, Any] = {
            'extract_images': options.pdf.extract_images,
            'extract_metadata': options.pdf.extract_metadata,
        }
        if options.pdf.enable_hierarchy:
            from kreuzberg import HierarchyConfig

            pdf_kwargs['hierarchy'] = HierarchyConfig(enabled=True)
        config_kwargs['pdf_options'] = PdfConfig(**pdf_kwargs)

    if options.language_detection:
        config_kwargs['language_detection'] = LanguageDetectionConfig(
            enabled=options.language_detection.enabled,
            min_confidence=options.language_detection.min_confidence,
            detect_multiple=options.language_detection.detect_multiple,
        )

    if options.token_reduction and options.token_reduction.mode:
        config_kwargs['token_reduction'] = TokenReductionConfig(
            mode=options.token_reduction.mode,
            preserve_important_words=options.token_reduction.preserve_important_words,
        )

    if options.pages:
        from kreuzberg import PageConfig

        page_kwargs: dict[str, Any] = {
            'extract_pages': options.pages.extract_pages,
            'insert_page_markers': options.pages.insert_page_markers,
        }
        if options.pages.marker_format:
            page_kwargs['marker_format'] = options.pages.marker_format
        config_kwargs['pages'] = PageConfig(**page_kwargs)

    if options.hierarchy and options.hierarchy.enabled:
        from kreuzberg import HierarchyConfig

        hierarchy_config = HierarchyConfig(
            enabled=True,
            k_clusters=options.hierarchy.k_clusters,
            include_bbox=options.hierarchy.include_bbox,
            ocr_coverage_threshold=options.hierarchy.ocr_coverage_threshold,
        )
        if 'pdf_options' not in config_kwargs:
            config_kwargs['pdf_options'] = PdfConfig(hierarchy=hierarchy_config)

    # New in kreuzberg 4.2.0: output format conversion
    if options.output_format:
        config_kwargs['output_format'] = options.output_format

    # New in kreuzberg 4.2.0: result format (unified vs element_based)
    if options.result_format:
        config_kwargs['result_format'] = options.result_format

    return ExtractionConfig(**config_kwargs) if config_kwargs else None


def _convert_result(kreuzberg_result: Any) -> ExtractionResult:
    """Convert Kreuzberg result to our ExtractionResult.

    Args:
        kreuzberg_result: Result from Kreuzberg extraction

    Returns:
        Our ExtractionResult dataclass
    """
    chunks = None
    if kreuzberg_result.chunks:
        chunks = []
        for i, chunk in enumerate(kreuzberg_result.chunks):
            # Kreuzberg returns chunks as dicts with 'content' key
            if isinstance(chunk, dict):
                chunk_content = chunk.get('content', '')
                metadata = chunk.get('metadata', {})
                start_char = metadata.get('byte_start', metadata.get('char_start'))
                end_char = metadata.get('byte_end', metadata.get('char_end'))
            else:
                chunk_content = getattr(chunk, 'content', str(chunk))
                start_char = getattr(chunk, 'start_char', None)
                end_char = getattr(chunk, 'end_char', None)
                metadata = getattr(chunk, 'metadata', None)

            chunks.append(
                TextChunk(
                    content=chunk_content,
                    index=i,
                    start_char=start_char,
                    end_char=end_char,
                    metadata=metadata if isinstance(metadata, dict) else None,
                )
            )

    return ExtractionResult(
        content=kreuzberg_result.content,
        chunks=chunks,
        mime_type=kreuzberg_result.mime_type,
        metadata=dict(kreuzberg_result.metadata) if kreuzberg_result.metadata else None,
        tables=getattr(kreuzberg_result, 'tables', None),
        detected_languages=getattr(kreuzberg_result, 'detected_languages', None),
    )


def extract_with_kreuzberg_sync(
    file_path: str | Path,
    options: ExtractionOptions | None = None,
) -> ExtractionResult:
    """Extract text from document using Kreuzberg (synchronous).

    This is the primary function for document extraction. It handles
    all 56+ formats supported by Kreuzberg.

    Args:
        file_path: Path to the document file
        options: Extraction options (chunking, OCR, etc.)

    Returns:
        ExtractionResult with content and optional chunks

    Raises:
        ImportError: If kreuzberg is not installed
        Exception: If extraction fails
    """
    if not is_kreuzberg_available():
        raise ImportError('kreuzberg is not installed. Install it with: pip install kreuzberg')

    from kreuzberg import extract_file_sync

    config = _build_extraction_config(options)

    file_str = str(file_path).lower()
    is_pptx_like = file_str.endswith(('.pptx', '.ppt', '.ppsx'))

    # PPSX uses PPTX MIME override - kreuzberg handles it if we lie about the MIME type
    effective_mime = options.mime_type if options else None
    if file_str.endswith('.ppsx') and not effective_mime:
        effective_mime = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        logger.debug('PPSX detected, using PPTX MIME type override for kreuzberg')

    try:
        result = extract_file_sync(
            file_path,
            mime_type=effective_mime,
            config=config,
        )
        return _convert_result(result)
    except Exception as e:
        # Kreuzberg PPTX parser fails on shapes without txBody (pictures, SmartArt, etc.)
        error_msg = str(e)
        if 'No txBody found' in error_msg and is_pptx_like:
            # For PPSX, use ZIP extraction (python-pptx also rejects PPSX)
            if file_str.endswith('.ppsx'):
                logger.warning(
                    f'Kreuzberg PPSX extraction failed ({error_msg}), using ZIP extraction fallback'
                )
                content = _extract_ppsx_from_zip(file_path)
                return ExtractionResult(
                    content=content,
                    mime_type='application/vnd.openxmlformats-officedocument.presentationml.slideshow',
                )
            # For PPTX, use python-pptx
            logger.warning(
                f'Kreuzberg PPTX extraction failed ({error_msg}), using python-pptx fallback'
            )
            content = _extract_pptx_with_python_pptx(file_path)
            return ExtractionResult(content=content, mime_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        logger.error(f'Kreuzberg extraction failed for {file_path}: {e}')
        raise


async def extract_with_kreuzberg(
    file_path: str | Path,
    options: ExtractionOptions | None = None,
) -> ExtractionResult:
    """Extract text from document using Kreuzberg (async).

    Async version that uses Kreuzberg's native async API for
    non-blocking document extraction.

    Args:
        file_path: Path to the document file
        options: Extraction options (chunking, OCR, etc.)

    Returns:
        ExtractionResult with content and optional chunks

    Raises:
        ImportError: If kreuzberg is not installed
        Exception: If extraction fails
    """
    if not is_kreuzberg_available():
        raise ImportError('kreuzberg is not installed. Install it with: pip install kreuzberg')

    from kreuzberg import extract_file

    config = _build_extraction_config(options)

    file_str = str(file_path).lower()
    is_pptx_like = file_str.endswith(('.pptx', '.ppt', '.ppsx'))

    # PPSX uses PPTX MIME override - kreuzberg handles it if we lie about the MIME type
    effective_mime = options.mime_type if options else None
    if file_str.endswith('.ppsx') and not effective_mime:
        effective_mime = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        logger.debug('PPSX detected, using PPTX MIME type override for kreuzberg')

    try:
        result = await extract_file(
            file_path,
            mime_type=effective_mime,
            config=config,
        )
        return _convert_result(result)
    except Exception as e:
        # Kreuzberg PPTX parser fails on shapes without txBody (pictures, SmartArt, etc.)
        error_msg = str(e)
        if 'No txBody found' in error_msg and is_pptx_like:
            # For PPSX, use ZIP extraction (python-pptx also rejects PPSX)
            if file_str.endswith('.ppsx'):
                logger.warning(
                    f'Kreuzberg PPSX extraction failed ({error_msg}), using ZIP extraction fallback'
                )
                content = _extract_ppsx_from_zip(file_path)
                return ExtractionResult(
                    content=content,
                    mime_type='application/vnd.openxmlformats-officedocument.presentationml.slideshow',
                )
            # For PPTX, use python-pptx
            logger.warning(
                f'Kreuzberg PPTX extraction failed ({error_msg}), using python-pptx fallback'
            )
            content = _extract_pptx_with_python_pptx(file_path)
            return ExtractionResult(content=content, mime_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        logger.error(f'Kreuzberg extraction failed for {file_path}: {e}')
        raise


def extract_bytes_with_kreuzberg_sync(
    data: bytes,
    mime_type: str,
    options: ExtractionOptions | None = None,
) -> ExtractionResult:
    """Extract text from bytes using Kreuzberg (synchronous).

    This function processes document data directly from memory without
    requiring a file on disk. Useful for S3/cloud storage workflows.

    Args:
        data: Document content as bytes
        mime_type: MIME type of the data (e.g., 'application/pdf')
        options: Extraction options (chunking, OCR, etc.)

    Returns:
        ExtractionResult with content and optional chunks

    Raises:
        ImportError: If kreuzberg is not installed
        Exception: If extraction fails
    """
    if not is_kreuzberg_available():
        raise ImportError('kreuzberg is not installed. Install it with: pip install kreuzberg')

    from kreuzberg import extract_bytes_sync

    config = _build_extraction_config(options)
    is_pptx_like = 'presentationml' in mime_type

    # PPSX (slideshow) not supported by kreuzberg or python-pptx, use ZIP extraction
    if 'slideshow' in mime_type:
        logger.info('PPSX MIME type detected, using ZIP extraction')
        content = _extract_ppsx_bytes_from_zip(data)
        return ExtractionResult(content=content, mime_type=mime_type)

    try:
        result = extract_bytes_sync(
            data,
            mime_type,
            config=config,
        )
        return _convert_result(result)
    except Exception as e:
        # Kreuzberg PPTX parser fails on shapes without txBody (pictures, SmartArt, etc.)
        # Fallback: use python-pptx which handles missing text gracefully
        error_msg = str(e)
        if 'No txBody found' in error_msg and is_pptx_like:
            logger.warning(
                f'Kreuzberg PPTX extraction failed ({error_msg}), using python-pptx fallback'
            )
            content = _extract_pptx_bytes_with_python_pptx(data)
            return ExtractionResult(content=content, mime_type=mime_type)

        logger.error(f'Kreuzberg bytes extraction failed for {mime_type}: {e}')
        raise


async def extract_bytes_with_kreuzberg(
    data: bytes,
    mime_type: str,
    options: ExtractionOptions | None = None,
) -> ExtractionResult:
    """Extract text from bytes using Kreuzberg (async).

    Async version that processes document data directly from memory
    without requiring a file on disk. Useful for S3/cloud storage workflows.

    Args:
        data: Document content as bytes
        mime_type: MIME type of the data (e.g., 'application/pdf')
        options: Extraction options (chunking, OCR, etc.)

    Returns:
        ExtractionResult with content and optional chunks

    Raises:
        ImportError: If kreuzberg is not installed
        Exception: If extraction fails
    """
    if not is_kreuzberg_available():
        raise ImportError('kreuzberg is not installed. Install it with: pip install kreuzberg')

    from kreuzberg import extract_bytes

    config = _build_extraction_config(options)
    is_pptx_like = 'presentationml' in mime_type

    # PPSX (slideshow) not supported by kreuzberg or python-pptx, use ZIP extraction
    if 'slideshow' in mime_type:
        logger.info('PPSX MIME type detected, using ZIP extraction')
        content = _extract_ppsx_bytes_from_zip(data)
        return ExtractionResult(content=content, mime_type=mime_type)

    try:
        result = await extract_bytes(
            data,
            mime_type,
            config=config,
        )
        return _convert_result(result)
    except Exception as e:
        # Kreuzberg PPTX parser fails on shapes without txBody (pictures, SmartArt, etc.)
        # Fallback: use python-pptx which handles missing text gracefully
        error_msg = str(e)
        if 'No txBody found' in error_msg and is_pptx_like:
            logger.warning(
                f'Kreuzberg PPTX extraction failed ({error_msg}), using python-pptx fallback'
            )
            content = _extract_pptx_bytes_with_python_pptx(data)
            return ExtractionResult(content=content, mime_type=mime_type)

        logger.error(f'Kreuzberg bytes extraction failed for {mime_type}: {e}')
        raise


async def batch_extract_with_kreuzberg(
    file_paths: list[str | Path],
    options: ExtractionOptions | None = None,
) -> list[ExtractionResult]:
    """Batch extract text from multiple documents (async).

    Uses Kreuzberg's batch processing for efficient multi-document extraction.

    Args:
        file_paths: List of paths to document files
        options: Extraction options applied to all files

    Returns:
        List of ExtractionResults in same order as inputs

    Raises:
        ImportError: If kreuzberg is not installed
        Exception: If batch extraction fails
    """
    if not is_kreuzberg_available():
        raise ImportError('kreuzberg is not installed. Install it with: pip install kreuzberg')

    from kreuzberg import batch_extract_files

    config = _build_extraction_config(options)

    try:
        results = await batch_extract_files(file_paths, config=config)
        return [_convert_result(r) for r in results]
    except Exception as e:
        logger.error(f'Kreuzberg batch extraction failed: {e}')
        raise


def create_chunking_options(
    chunk_token_size: int = 1200,
    chunk_overlap_token_size: int = 100,
    preset: str | None = 'semantic',
) -> ChunkingOptions:
    """Create ChunkingOptions from YAR's token-based settings.

    This is a convenience function that converts YAR's token-based
    chunking parameters to Kreuzberg's character-based parameters.

    Args:
        chunk_token_size: Maximum tokens per chunk (default: 1200)
        chunk_overlap_token_size: Overlapping tokens between chunks (default: 100)
        preset: Chunking preset - 'semantic', 'recursive', or None

    Returns:
        ChunkingOptions configured for one-pass extraction+chunking
    """
    max_chars = tokens_to_chars(chunk_token_size)
    max_overlap = tokens_to_chars(chunk_overlap_token_size)

    return ChunkingOptions(
        enabled=True,
        max_chars=max_chars,
        max_overlap=max_overlap,
        preset=preset,
    )


def create_markdown_options(
    chunking: ChunkingOptions | None = None,
    ocr: OcrOptions | None = None,
) -> ExtractionOptions:
    """Create options for extracting document content as Markdown.

    New in kreuzberg 4.2.0: Any document format can be converted to Markdown
    during extraction, preserving structure like headings, lists, and tables.

    Args:
        chunking: Optional chunking configuration
        ocr: Optional OCR configuration for scanned documents

    Returns:
        ExtractionOptions configured for Markdown output

    Example:
        >>> options = create_markdown_options()
        >>> result = extract_with_kreuzberg_sync("report.pdf", options)
        >>> print(result.content)  # Markdown-formatted content
    """
    return ExtractionOptions(
        output_format='markdown',
        chunking=chunking,
        ocr=ocr,
    )


def create_html_options(
    chunking: ChunkingOptions | None = None,
    ocr: OcrOptions | None = None,
) -> ExtractionOptions:
    """Create options for extracting document content as HTML.

    New in kreuzberg 4.2.0: Any document format can be converted to HTML
    during extraction, preserving rich formatting and structure.

    Args:
        chunking: Optional chunking configuration
        ocr: Optional OCR configuration for scanned documents

    Returns:
        ExtractionOptions configured for HTML output
    """
    return ExtractionOptions(
        output_format='html',
        chunking=chunking,
        ocr=ocr,
    )


def extract_and_chunk_sync(
    file_path: str | Path,
    chunk_token_size: int = 1200,
    chunk_overlap_token_size: int = 100,
    chunking_preset: str | None = 'semantic',
    ocr_options: OcrOptions | None = None,
) -> ExtractionResult:
    """One-pass extraction and chunking from a document file (synchronous).

    This function performs document extraction and semantic chunking in a single
    pass, preserving document structure for better chunk boundaries. This is more
    efficient and produces better results than extracting text first and then
    re-chunking it.

    Args:
        file_path: Path to the document file
        chunk_token_size: Maximum tokens per chunk (default: 1200)
        chunk_overlap_token_size: Overlapping tokens between chunks (default: 100)
        chunking_preset: Chunking preset - 'semantic', 'recursive', or None
        ocr_options: Optional OCR configuration for scanned documents/images

    Returns:
        ExtractionResult with content and chunks

    Raises:
        ImportError: If kreuzberg is not installed
        Exception: If extraction fails
    """
    chunking_options = create_chunking_options(
        chunk_token_size=chunk_token_size,
        chunk_overlap_token_size=chunk_overlap_token_size,
        preset=chunking_preset,
    )

    options = ExtractionOptions(
        chunking=chunking_options,
        ocr=ocr_options,
    )

    return extract_with_kreuzberg_sync(file_path, options)


async def extract_and_chunk(
    file_path: str | Path,
    chunk_token_size: int = 1200,
    chunk_overlap_token_size: int = 100,
    chunking_preset: str | None = 'semantic',
    ocr_options: OcrOptions | None = None,
) -> ExtractionResult:
    """One-pass extraction and chunking from a document file (async).

    This function performs document extraction and semantic chunking in a single
    pass, preserving document structure for better chunk boundaries. This is more
    efficient and produces better results than extracting text first and then
    re-chunking it.

    Args:
        file_path: Path to the document file
        chunk_token_size: Maximum tokens per chunk (default: 1200)
        chunk_overlap_token_size: Overlapping tokens between chunks (default: 100)
        chunking_preset: Chunking preset - 'semantic', 'recursive', or None
        ocr_options: Optional OCR configuration for scanned documents/images

    Returns:
        ExtractionResult with content and chunks

    Raises:
        ImportError: If kreuzberg is not installed
        Exception: If extraction fails
    """
    chunking_options = create_chunking_options(
        chunk_token_size=chunk_token_size,
        chunk_overlap_token_size=chunk_overlap_token_size,
        preset=chunking_preset,
    )

    options = ExtractionOptions(
        chunking=chunking_options,
        ocr=ocr_options,
    )

    return await extract_with_kreuzberg(file_path, options)


def chunks_to_yar_format(result: ExtractionResult) -> list[dict[str, Any]]:
    """Convert ExtractionResult chunks to YAR's expected chunk format.

    This converts Kreuzberg's TextChunk objects to the dictionary format
    expected by YAR's document processing pipeline.

    Args:
        result: ExtractionResult from extraction (must have chunks)

    Returns:
        List of chunk dictionaries compatible with YAR's chunking_func output
    """
    if not result.chunks:
        # Fallback: return whole content as single chunk
        return [
            {
                'tokens': len(result.content) // 4,
                'content': result.content.strip(),
                'chunk_order_index': 0,
                'char_start': 0,
                'char_end': len(result.content),
            }
        ]

    chunks = []
    for chunk in result.chunks:
        # Estimate token count (~4 chars per token)
        tokens = len(chunk.content) // 4

        chunks.append(
            {
                'tokens': tokens,
                'content': chunk.content.strip(),
                'chunk_order_index': chunk.index,
                'char_start': chunk.start_char if chunk.start_char is not None else 0,
                'char_end': chunk.end_char if chunk.end_char is not None else len(chunk.content),
            }
        )

    return chunks


def get_supported_formats() -> list[str]:
    """Get list of file extensions supported by Kreuzberg.

    Returns:
        List of supported extensions (e.g., ['.pdf', '.docx', ...])
    """
    # Kreuzberg supports 56+ formats. Here are the most common ones.
    # Full list available at kreuzberg.dev/formats
    return [
        # Documents
        '.pdf',
        '.docx',
        '.doc',
        '.odt',
        '.rtf',
        '.txt',
        '.md',
        '.markdown',
        # Presentations
        '.pptx',
        '.ppsx',  # PowerPoint Show (same structure as PPTX)
        '.ppt',
        '.odp',
        # Spreadsheets
        '.xlsx',
        '.xls',
        '.ods',
        '.csv',
        # E-books
        '.epub',
        '.mobi',
        # Web
        '.html',
        '.htm',
        '.xml',
        '.json',
        # Images (for OCR)
        '.png',
        '.jpg',
        '.jpeg',
        '.gif',
        '.bmp',
        '.tiff',
        '.webp',
        # Archives (extracts contained docs)
        '.zip',
        '.tar',
        '.gz',
        # Email
        '.eml',
        '.msg',
    ]
